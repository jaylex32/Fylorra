"""
Fylorra - Library Index
Local SQLite index for file metadata, extracted text, and AI summaries.
"""

from __future__ import annotations

import hashlib
import json
import os
import sqlite3
import time
import tempfile
import logging
import warnings
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Optional


# Quiet noisy PDF parsers in console logs.
logging.getLogger("pypdf").setLevel(logging.ERROR)
logging.getLogger("fitz").setLevel(logging.ERROR)
warnings.filterwarnings("ignore", module="pypdf")
warnings.filterwarnings("ignore", message="Ignoring wrong pointing object*")
warnings.filterwarnings("ignore", message="startxref on same line as offset*")
warnings.filterwarnings("ignore", message="Unexpected escaped string*")


@dataclass(frozen=True)
class LibraryItem:
    path: str
    name: str
    ext: str
    size: int
    mtime: float
    sha256: Optional[str]
    extracted_text: Optional[str]
    ai_summary: Optional[str]
    ai_tags_json: Optional[str]


class LibraryIndex:
    """
    SQLite index with optional FTS5 for fast local search.

    Notes:
    - Uses only stdlib (sqlite3) and is safe to ship.
    - FTS5 availability depends on the Python sqlite build; we fall back to LIKE.
    """

    def __init__(self, db_path: Optional[Path] = None):
        if db_path is None:
            app_data = Path.home() / ".fylorra"
            app_data.mkdir(exist_ok=True)
            db_path = app_data / "library_index.db"
        self.db_path = Path(db_path)
        self._fts_enabled = False
        self._init_database()

    @property
    def fts_enabled(self) -> bool:
        return self._fts_enabled

    def _init_database(self) -> None:
        self.db_path.parent.mkdir(parents=True, exist_ok=True)
        with sqlite3.connect(self.db_path) as conn:
            cur = conn.cursor()
            cur.execute(
                """
                CREATE TABLE IF NOT EXISTS files (
                    path TEXT PRIMARY KEY,
                    name TEXT NOT NULL,
                    ext TEXT NOT NULL,
                    size INTEGER NOT NULL,
                    mtime REAL NOT NULL,
                    sha256 TEXT,
                    extracted_text TEXT,
                    ai_summary TEXT,
                    ai_tags_json TEXT,
                    updated_at REAL NOT NULL
                )
                """
            )
            cur.execute("CREATE INDEX IF NOT EXISTS idx_files_mtime ON files(mtime DESC)")
            cur.execute("CREATE INDEX IF NOT EXISTS idx_files_ext ON files(ext)")

            # Try enabling FTS5 (optional)
            try:
                cur.execute(
                    """
                    CREATE VIRTUAL TABLE IF NOT EXISTS files_fts
                    USING fts5(path, name, ext, extracted_text, ai_summary, ai_tags_json)
                    """
                )
                self._fts_enabled = True
            except sqlite3.OperationalError:
                self._fts_enabled = False

            conn.commit()

    @staticmethod
    def _walk_files(
        folder: Path,
        *,
        include_subfolders: bool,
        include_hidden: bool,
        allow_exts: Optional[set[str]],
        max_files: Optional[int],
    ) -> list[Path]:
        def is_hidden_name(name: str) -> bool:
            return bool(name) and name.startswith(".")

        # Common internal/output folders we should never index by default.
        # This prevents accidental "index loops" on conversion output trees.
        exclude_dir_names = {
            "__pycache__",
            ".git",
            ".hg",
            ".svn",
            ".fylorra",
            "Converted_Media",
            "Converted_Images",
            "Converted_Office",
            "Converted_PDF",
            "Converted_Archives",
            "Split_Pages",
            "Split_Chunks",
            "Split_By_Bookmarks",
            "MP3 Music",
        }

        folder = Path(folder)
        paths: list[Path] = []

        def allow_file(p: Path) -> bool:
            if not include_hidden and (is_hidden_name(p.name) or LibraryIndex._is_hidden(p)):
                return False
            if allow_exts is not None and p.suffix.lower() not in allow_exts:
                return False
            return True

        if not include_subfolders:
            try:
                for entry in os.scandir(folder):
                    if not entry.is_file():
                        continue
                    p = Path(entry.path)
                    if allow_file(p):
                        paths.append(p)
                        if max_files and len(paths) >= max_files:
                            break
            except Exception:
                return paths
            return paths

        for root, dirs, files in os.walk(folder):
            try:
                root_path = Path(root)
                dirs[:] = [d for d in dirs if d not in exclude_dir_names]
                if not include_hidden:
                    dirs[:] = [d for d in dirs if not is_hidden_name(d)]
            except Exception:
                pass

            for name in files:
                p = Path(root) / name
                if allow_file(p):
                    paths.append(p)
                    if max_files and len(paths) >= max_files:
                        return paths
        return paths

    @staticmethod
    def _fts_safe_query(raw: str) -> str:
        """
        Make a query safe for FTS MATCH. This is used as a fallback when the raw query
        triggers FTS syntax errors (unbalanced quotes, stray operators, etc).
        """
        q = (raw or "").strip()
        if not q:
            return q

        q = q.replace("\\", " ").replace("/", " ")
        q = re.sub(r"[\x00-\x1f]", " ", q)
        # Remove common FTS operators that users/models may accidentally emit.
        q = re.sub(r"\b(NEAR|NOT)\b", " ", q, flags=re.IGNORECASE)
        q = q.replace(":", " ")

        tokens: list[str] = []
        for t in q.split():
            t2 = re.sub(r"[^\w\-.]+", "", t, flags=re.UNICODE).strip("._-")
            if len(t2) < 2:
                continue
            if t2.upper() in {"OR", "AND"}:
                continue
            tokens.append(t2)

        if not tokens:
            return q
        if len(tokens) == 1:
            return f"\"{tokens[0].replace('\"', '')}\"*"
        # OR is much more forgiving for natural language searches.
        tokens = tokens[:10]
        return " OR ".join(f"\"{t.replace('\"', '')}\"*" for t in tokens)

    def _upsert(self, item: LibraryItem) -> None:
        now = time.time()
        with sqlite3.connect(self.db_path) as conn:
            cur = conn.cursor()
            self._upsert_with_cursor(cur, item, now=now)
            conn.commit()

    def _upsert_with_cursor(self, cur: sqlite3.Cursor, item: LibraryItem, *, now: float) -> None:
        cur.execute(
            """
            INSERT INTO files(path, name, ext, size, mtime, sha256, extracted_text, ai_summary, ai_tags_json, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(path) DO UPDATE SET
                name=excluded.name,
                ext=excluded.ext,
                size=excluded.size,
                mtime=excluded.mtime,
                sha256=excluded.sha256,
                extracted_text=excluded.extracted_text,
                ai_summary=excluded.ai_summary,
                ai_tags_json=excluded.ai_tags_json,
                updated_at=excluded.updated_at
            """,
            (
                item.path,
                item.name,
                item.ext,
                int(item.size),
                float(item.mtime),
                item.sha256,
                item.extracted_text,
                item.ai_summary,
                item.ai_tags_json,
                now,
            ),
        )
        if self._fts_enabled:
            row = (
                item.path,
                item.name,
                item.ext,
                item.extracted_text or "",
                item.ai_summary or "",
                item.ai_tags_json or "",
            )
            cur.execute("DELETE FROM files_fts WHERE path = ?", (item.path,))
            cur.execute(
                "INSERT INTO files_fts(path, name, ext, extracted_text, ai_summary, ai_tags_json) VALUES (?, ?, ?, ?, ?, ?)",
                row,
            )

    @staticmethod
    def _sha256(path: Path, *, max_bytes: int = 20 * 1024 * 1024) -> Optional[str]:
        try:
            size = path.stat().st_size
            if size > max_bytes:
                return None
            h = hashlib.sha256()
            with path.open("rb") as f:
                for chunk in iter(lambda: f.read(1024 * 1024), b""):
                    h.update(chunk)
            return h.hexdigest()
        except Exception:
            return None

    @staticmethod
    def _safe_stat(path: Path) -> Optional[tuple[int, float]]:
        try:
            st = path.stat()
            return int(st.st_size), float(st.st_mtime)
        except Exception:
            return None

    @staticmethod
    def _is_hidden(path: Path) -> bool:
        name = path.name
        if name.startswith("."):
            return True
        try:
            # Windows hidden attribute
            return bool(os.stat(path).st_file_attributes & 2)  # type: ignore[attr-defined]
        except Exception:
            return False

    def index_folder(
        self,
        folder: Path,
        *,
        include_subfolders: bool = True,
        include_hidden: bool = False,
        ai_manager=None,
        ai_summarize: bool = False,
        max_ai_files: int = 200,
        extract_images: bool = False,
        compute_hashes: bool = False,
        include_extensions: Optional[list[str]] = None,
        ocr_scanned_pdfs: bool = False,
        ocr_pdf_pages: int = 1,
        max_pdf_ocr_files: int = 60,
        progress_cb: Optional[Callable[[str, float], None]] = None,
        max_files: Optional[int] = None,
    ) -> int:
        """
        Index files in a folder. Extracts lightweight text and optional AI summary.

        progress_cb(message, progress_0_to_1)
        """
        folder = Path(folder)
        if not folder.exists() or not folder.is_dir():
            return 0

        allow_exts = None
        if include_extensions:
            allow_exts = {("." + str(e).lower().lstrip(".")).strip() for e in include_extensions if str(e).strip()}

        paths = self._walk_files(
            folder,
            include_subfolders=include_subfolders,
            include_hidden=include_hidden,
            allow_exts=allow_exts,
            max_files=max_files,
        )

        total = len(paths)
        if total == 0:
            return 0

        ai_done = 0
        ocr_done = 0
        unchanged = 0
        now = time.time()

        # Single connection for speed (transactions)
        with sqlite3.connect(self.db_path) as conn:
            cur = conn.cursor()
            cur.execute("PRAGMA journal_mode=WAL")
            cur.execute("PRAGMA synchronous=NORMAL")
            cur.execute("PRAGMA temp_store=MEMORY")

            for idx, p in enumerate(paths, start=1):
                rel = None
                try:
                    rel = str(p.relative_to(folder))
                except Exception:
                    rel = p.name
                if progress_cb:
                    progress_cb(f"Indexing {idx}/{total}: {rel}", (idx - 1) / total)

                stat = self._safe_stat(p)
                if not stat:
                    continue
                size, mtime = stat
                ext = p.suffix.lower()

                # Skip unchanged files (major speed-up for repeat indexing).
                try:
                    row = cur.execute(
                        "SELECT size, mtime, sha256, extracted_text, ai_summary FROM files WHERE path = ?",
                        (str(p),),
                    ).fetchone()
                    if row and int(row[0]) == int(size) and float(row[1]) == float(mtime):
                        needs_refresh = False

                        # If the user enabled new enrichment options, we must re-process
                        # unchanged files that are missing those fields.
                        if compute_hashes:
                            try:
                                if row[2] is None and int(size) <= 20 * 1024 * 1024:
                                    needs_refresh = True
                            except Exception:
                                pass

                        if extract_images and ai_manager and getattr(ai_manager, "is_ready", False):
                            try:
                                if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}:
                                    if not (row[3] or "").strip():
                                        needs_refresh = True
                            except Exception:
                                pass

                        if ai_summarize and ai_manager and getattr(ai_manager, "is_ready", False):
                            try:
                                # If we have extracted_text but no summary, allow summarization.
                                if (row[3] or "").strip() and not (row[4] or "").strip() and ai_done < int(max_ai_files):
                                    needs_refresh = True
                            except Exception:
                                pass

                        if (
                            ext == ".pdf"
                            and ocr_scanned_pdfs
                            and ai_manager
                            and getattr(ai_manager, "is_ready", False)
                            and ocr_done < int(max_pdf_ocr_files)
                        ):
                            try:
                                if not (row[3] or "").strip():
                                    needs_refresh = True
                            except Exception:
                                pass

                        if not needs_refresh:
                            unchanged += 1
                            if progress_cb:
                                progress_cb(f"Skipped unchanged {idx}/{total}: {rel}", idx / total)
                            continue
                except Exception:
                    pass

                extracted_text = None
                ai_summary = None
                ai_tags_json = None

                # Lightweight extraction:
                # - For images, reuse existing AIManager helper (describes visible text/content).
                # - For PDFs and text-like files, extract local text for search (no AI required).
                if extract_images and ai_manager and getattr(ai_manager, "is_ready", False):
                    try:
                        if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}:
                            # Prefer a semantic vision caption + tags for image search.
                            # Fallback to older helper names if present.
                            cap = None
                            tags: list[str] = []
                            fn_cap = getattr(ai_manager, "extract_image_caption_tags", None)
                            if callable(fn_cap):
                                out = fn_cap(p)
                                if isinstance(out, dict):
                                    cap = (out.get("caption") or "").strip() or None
                                    raw_tags = out.get("tags") or []
                                    if isinstance(raw_tags, list):
                                        tags = [str(t) for t in raw_tags if str(t).strip()]

                            if not cap:
                                fn = getattr(ai_manager, "extract_text_description", None) or getattr(ai_manager, "extract_text_content", None)
                                if callable(fn):
                                    cap = (fn(p) or "").strip() or None

                            # Store as searchable text (caption + tags) so filename-only images still match.
                            if cap or tags:
                                extracted_text = " ".join([cap or "", " ".join(tags)]).strip() or None
                                if tags:
                                    try:
                                        ai_tags_json = json.dumps({"image_tags": tags}, ensure_ascii=False)
                                    except Exception:
                                        ai_tags_json = None
                                ai_summary = (cap or extracted_text or "")[:500] if (cap or extracted_text) else None
                    except Exception:
                        pass

                if extracted_text is None:
                    extracted_text = self._extract_local_text(p, ext=ext, max_chars=30_000)
                    if extracted_text:
                        if ai_summarize and ai_manager and getattr(ai_manager, "is_ready", False) and ai_done < int(max_ai_files):
                            try:
                                ai_summary = self._ai_summarize_text(ai_manager, p.name, extracted_text)
                                ai_done += 1
                            except Exception:
                                ai_summary = None
                        if not ai_summary:
                            ai_summary = extracted_text[:500]
                    elif (
                        ext == ".pdf"
                        and ocr_scanned_pdfs
                        and ai_manager
                        and getattr(ai_manager, "is_ready", False)
                        and ocr_done < int(max_pdf_ocr_files)
                    ):
                        # Scanned PDF (no extracted text): render + vision OCR for the first N pages.
                        try:
                            extracted_text = self._ocr_pdf_with_ai(ai_manager, p, max_pages=max(1, int(ocr_pdf_pages)))
                            if extracted_text:
                                ai_summary = extracted_text[:500]
                                ocr_done += 1
                        except Exception:
                            pass

                sha256 = None
                # Hashing is optional (useful for dedupe, but slow for large media/photo libraries).
                if compute_hashes and size <= 20 * 1024 * 1024:
                    sha256 = self._sha256(p)

                item = LibraryItem(
                    path=str(p),
                    name=p.name,
                    ext=ext,
                    size=size,
                    mtime=mtime,
                    sha256=sha256,
                    extracted_text=extracted_text,
                    ai_summary=ai_summary,
                    ai_tags_json=ai_tags_json,
                )
                self._upsert_with_cursor(cur, item, now=now)
                if idx % 200 == 0:
                    conn.commit()
                if progress_cb:
                    progress_cb(f"Indexed {idx}/{total}: {rel}", idx / total)

            conn.commit()

        if progress_cb:
            msg = f"Index complete. Indexed {total - unchanged}/{total} (skipped unchanged {unchanged})."
            if ocr_done:
                msg += f" OCR'd {ocr_done} scanned PDF(s)."
            progress_cb(msg, 1.0)
        return total

    @staticmethod
    def _ocr_pdf_with_ai(ai_manager, pdf_path: Path, *, max_pages: int = 1) -> Optional[str]:
        """
        OCR scanned PDFs by rendering pages and using the local vision model.
        Requires PyMuPDF (fitz). No external OCR apps needed.
        """
        try:
            import fitz  # type: ignore
        except Exception:
            return None

        max_pages = max(1, int(max_pages))
        text_parts: list[str] = []
        with tempfile.TemporaryDirectory() as td:
            doc = fitz.open(str(pdf_path))
            try:
                page_count = min(doc.page_count, max_pages)
                for i in range(page_count):
                    page = doc.load_page(i)
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                    img_path = Path(td) / f"page_{i+1}.png"
                    pix.save(str(img_path))
                    t = ai_manager.extract_text_content(img_path)
                    if t:
                        text_parts.append(str(t))
            finally:
                try:
                    doc.close()
                except Exception:
                    pass

        out = "\n".join(text_parts).strip()
        return out[:30_000] if out else None

    @staticmethod
    def _ai_summarize_text(ai_manager, filename: str, text: str) -> Optional[str]:
        """
        Create a short searchable summary (<= ~2 lines) from extracted text.
        """
        text = (text or "").strip()
        if not text:
            return None
        snippet = text[:4000]
        prompt = (
            "Summarize this file content for LOCAL SEARCH.\n"
            "Rules:\n"
            "- 1-2 lines max.\n"
            "- Include strong keywords: vendor, doc type, identifiers, dates.\n"
            "- No markdown.\n\n"
            f"Filename: {filename}\n\n"
            f"Content:\n{snippet}\n"
        )
        resp = ai_manager.model.create_chat_completion(
            messages=[{"role": "user", "content": [{"type": "text", "text": prompt}]}],
            temperature=0.2,
            max_tokens=120,
        )
        out = resp["choices"][0]["message"]["content"]
        out = (out or "").strip()
        return out[:600] if out else None

    @staticmethod
    def _extract_local_text(file_path: Path, *, ext: str, max_chars: int) -> Optional[str]:
        """
        Best-effort local text extraction for indexing.
        Keeps runtime optional: no hard dependency on external libs.
        """
        try:
            if ext in {".txt", ".md", ".log", ".csv", ".json", ".yaml", ".yml"}:
                # Avoid huge reads; enough for search.
                with file_path.open("rb") as f:
                    data = f.read(512 * 1024)
                text = data.decode("utf-8", errors="ignore").strip()
                return text[:max_chars] if text else None

            if ext == ".pdf":
                # Try PyMuPDF first (best extraction), then fall back to pypdf.
                try:
                    import fitz  # type: ignore

                    doc = fitz.open(str(file_path))
                    parts: list[str] = []
                    for i in range(min(doc.page_count, 10)):
                        parts.append(doc.load_page(i).get_text("text"))
                    doc.close()
                    text = "\n".join(parts).strip()
                    return text[:max_chars] if text else None
                except Exception:
                    pass

                try:
                    import warnings

                    logging.getLogger("pypdf").setLevel(logging.ERROR)
                    warnings.filterwarnings("ignore", module="pypdf")
                    from pypdf import PdfReader  # type: ignore

                    reader = PdfReader(str(file_path))
                    parts = []
                    for page in reader.pages[:10]:
                        parts.append(page.extract_text() or "")
                    text = "\n".join(parts).strip()
                    return text[:max_chars] if text else None
                except Exception:
                    return None

            if ext == ".docx":
                try:
                    import docx  # type: ignore

                    d = docx.Document(str(file_path))
                    text = "\n".join([p.text for p in d.paragraphs if p.text]).strip()
                    return text[:max_chars] if text else None
                except Exception:
                    pass

                # Fallback: LibreOffice convert to PDF then extract.
                return LibraryIndex._extract_text_via_libreoffice_pdf(file_path, max_chars=max_chars)

            if ext == ".pptx":
                try:
                    from pptx import Presentation  # type: ignore

                    prs = Presentation(str(file_path))
                    parts: list[str] = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            t = getattr(shape, "text", None)
                            if t:
                                parts.append(str(t))
                    text = "\n".join(parts).strip()
                    return text[:max_chars] if text else None
                except Exception:
                    pass

                return LibraryIndex._extract_text_via_libreoffice_pdf(file_path, max_chars=max_chars)

            if ext == ".xlsx":
                try:
                    import openpyxl  # type: ignore

                    wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
                    parts: list[str] = []
                    for sheet in wb.worksheets[:5]:
                        parts.append(sheet.title)
                        row_count = 0
                        for row in sheet.iter_rows(values_only=True):
                            row_count += 1
                            if row_count > 50:
                                break
                            vals = [str(v) for v in row if v is not None and str(v).strip()]
                            if vals:
                                parts.append(" | ".join(vals))
                    try:
                        wb.close()
                    except Exception:
                        pass
                    text = "\n".join(parts).strip()
                    return text[:max_chars] if text else None
                except Exception:
                    pass

                return LibraryIndex._extract_text_via_libreoffice_pdf(file_path, max_chars=max_chars)

        except Exception:
            return None

        return None

    @staticmethod
    def _extract_text_via_libreoffice_pdf(file_path: Path, *, max_chars: int) -> Optional[str]:
        """
        Best-effort Office→PDF conversion via LibreOffice (if installed), then extract text from the PDF.
        """
        try:
            from core.lo_converter import LibreOfficeConverter
        except Exception:
            return None

        try:
            conv = LibreOfficeConverter()
            if not conv.is_available():
                return None

            with tempfile.TemporaryDirectory(prefix="fylorra_index_") as td:
                out_dir = Path(td)
                pdf_path = conv.convert_to_pdf(file_path, out_dir=out_dir)
                if not pdf_path or not pdf_path.exists():
                    return None

                return LibraryIndex._extract_local_text(pdf_path, ext=".pdf", max_chars=max_chars)
        except Exception:
            return None

    def search(self, query: str, *, limit: int = 50, folder: Optional[Path] = None) -> list[LibraryItem]:
        q = (query or "").strip()
        if not q:
            return []

        folder_prefix = None
        folder_prefix_slash = None
        # Params:
        # 0 prefix (native), 1 prefix\%, 2 prefix/%, 3 prefix_slash/%
        folder_params: tuple[str, str, str, str] | None = None
        if folder:
            try:
                folder_prefix = str(Path(folder).resolve())
                folder_prefix = folder_prefix.rstrip("\\/")  # normalize
                folder_prefix_slash = folder_prefix.replace("\\", "/")
                folder_params = (
                    folder_prefix,
                    folder_prefix + "\\%",
                    folder_prefix + "/%",
                    folder_prefix_slash + "/%",
                )
            except Exception:
                folder_prefix = None
                folder_prefix_slash = None
                folder_params = None

        def _like_tokens(raw: str) -> list[str]:
            raw = (raw or "").strip()
            if not raw:
                return []
            tokens: list[str] = []
            for t in re.split(r"\s+", raw.replace("\\", " ").replace("/", " ")):
                t = re.sub(r"[^\w\-.]+", "", t, flags=re.UNICODE).strip()
                if len(t) < 3:
                    continue
                tokens.append(t)
                # Cheap singularization (invoice(s), errors, etc) for better LIKE hits.
                if len(t) > 4 and t.lower().endswith("s"):
                    tokens.append(t[:-1])
            # De-dupe while preserving order.
            seen: set[str] = set()
            out: list[str] = []
            for t in tokens:
                k = t.lower()
                if k in seen:
                    continue
                seen.add(k)
                out.append(t)
                if len(out) >= 8:
                    break
            return out

        def _run_like(cur: sqlite3.Cursor) -> list[sqlite3.Row]:
            tokens = _like_tokens(q)
            if not tokens:
                tokens = [q]
            clauses: list[str] = []
            params: list[str] = []
            for t in tokens:
                like = f"%{t}%"
                clauses.append("(name LIKE ? OR extracted_text LIKE ? OR ai_summary LIKE ? OR ai_tags_json LIKE ?)")
                params.extend([like, like, like, like])
            where = " OR ".join(clauses) if clauses else "1=0"
            folder_sql = ""
            if folder_params:
                # Some historical rows may store paths with mixed slashes (e.g. C:/Users/...\\file.jpg).
                # Use a normalized comparison as well to ensure folder scoping always works.
                folder_sql = (
                    " AND ("
                    "path = ? OR path LIKE ? OR path LIKE ? "
                    "OR replace(path,'\\\\','/') LIKE ?"
                    ")"
                )
            sql = (
                "SELECT * FROM files "
                f"WHERE {where}{folder_sql} "
                "ORDER BY mtime DESC "
                "LIMIT ?"
            )
            if folder_params:
                cur.execute(sql, (*params, *folder_params, int(limit)))
            else:
                cur.execute(sql, (*params, int(limit)))
            return cur.fetchall()

        with sqlite3.connect(self.db_path) as conn:
            conn.row_factory = sqlite3.Row
            cur = conn.cursor()
            rows: list[sqlite3.Row] = []
            if self._fts_enabled:
                ran_fts = False
                for attempt in range(2):
                    try_q = q if attempt == 0 else self._fts_safe_query(q)
                    try:
                        if folder_params:
                            cur.execute(
                                """
                                SELECT f.*, bm25(files_fts) AS rank
                                FROM files_fts
                                JOIN files f ON f.path = files_fts.path
                                WHERE files_fts MATCH ?
                                  AND (
                                    f.path = ? OR f.path LIKE ? OR f.path LIKE ?
                                    OR replace(f.path,'\\','/') LIKE ?
                                  )
                                ORDER BY rank ASC, f.mtime DESC
                                LIMIT ?
                                """,
                                (try_q, *folder_params, limit),
                            )
                        else:
                            cur.execute(
                                """
                                SELECT f.*, bm25(files_fts) AS rank
                                FROM files_fts
                                JOIN files f ON f.path = files_fts.path
                                WHERE files_fts MATCH ?
                                ORDER BY rank ASC, f.mtime DESC
                                LIMIT ?
                                """,
                                (try_q, limit),
                            )
                        ran_fts = True
                        break
                    except sqlite3.OperationalError:
                        continue
                if not ran_fts:
                    rows = _run_like(cur)
                else:
                    rows = cur.fetchall()
                    # FTS can be too strict (no stemming). If it yields nothing, fall back to LIKE tokens.
                    if not rows:
                        rows = _run_like(cur)
            else:
                rows = _run_like(cur)

        # Rank results to prioritize "obvious" matches:
        # - filename/path hits > extracted text hits
        # - common document types (pdf/doc/xls) > everything else
        try:
            tokens = _like_tokens(q)
            if not tokens:
                tokens = [q]

            ext_bonus = {
                ".pdf": 10,
                ".doc": 9,
                ".docx": 9,
                ".xls": 9,
                ".xlsx": 9,
                ".ppt": 8,
                ".pptx": 8,
                ".csv": 6,
                ".txt": 5,
                ".md": 3,
                ".png": 2,
                ".jpg": 2,
                ".jpeg": 2,
                ".webp": 2,
            }

            scored: list[tuple[int, float, float, sqlite3.Row]] = []
            for row in rows:
                try:
                    name = (row["name"] or "")
                    path = (row["path"] or "")
                    ext = (row["ext"] or "")
                    extracted = (row["extracted_text"] or "")
                    summary = (row["ai_summary"] or "")
                    mtime = float(row["mtime"] or 0.0)
                except Exception:
                    continue

                nlow = name.lower()
                plow = path.lower()
                elow = ext.lower()
                slow = summary.lower()
                xlow = extracted.lower()

                score = int(ext_bonus.get(elow, 0))
                for t in tokens:
                    tl = t.lower()
                    if not tl or len(tl) < 3:
                        continue
                    if tl in nlow:
                        score += 50
                    if tl in plow:
                        score += 12
                    if tl in slow:
                        score += 8
                    if tl in xlow:
                        score += 4

                # If FTS provided a bm25 rank, prefer lower ranks.
                rank = 9999.0
                try:
                    if "rank" in row.keys():
                        rank = float(row["rank"])
                except Exception:
                    rank = 9999.0

                scored.append((score, mtime, rank, row))

            # Highest score first, then newer, then better bm25 rank.
            scored.sort(key=lambda x: (-x[0], -x[1], x[2]))
            rows = [r for _s, _m, _rk, r in scored][: int(limit)]
        except Exception:
            pass

        return [
            LibraryItem(
                path=row["path"],
                name=row["name"],
                ext=row["ext"],
                size=int(row["size"]),
                mtime=float(row["mtime"]),
                sha256=row["sha256"],
                extracted_text=row["extracted_text"],
                ai_summary=row["ai_summary"],
                ai_tags_json=row["ai_tags_json"],
            )
            for row in rows
        ]

    def get(self, path: Path) -> Optional[LibraryItem]:
        with sqlite3.connect(self.db_path) as conn:
            conn.row_factory = sqlite3.Row
            cur = conn.cursor()
            cur.execute("SELECT * FROM files WHERE path = ?", (str(path),))
            row = cur.fetchone()
        if not row:
            return None
        return LibraryItem(
            path=row["path"],
            name=row["name"],
            ext=row["ext"],
            size=int(row["size"]),
            mtime=float(row["mtime"]),
            sha256=row["sha256"],
            extracted_text=row["extracted_text"],
            ai_summary=row["ai_summary"],
            ai_tags_json=row["ai_tags_json"],
        )
