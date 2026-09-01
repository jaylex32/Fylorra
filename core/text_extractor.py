"""
Fylorra - Text Extractor (OCR + local parsing)
Offline-first extraction for PDFs, images, and office files.
"""

from __future__ import annotations

from dataclasses import dataclass
import base64
import importlib.util
import io
import re
from pathlib import Path
from typing import Callable, Optional


TEXT_EXTS = {".txt", ".md", ".log", ".csv", ".json", ".xml", ".yaml", ".yml"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".gif", ".webp"}
DOC_EXTS = {".docx"}
PPT_EXTS = {".pptx"}
SHEET_EXTS = {".xlsx"}


@dataclass(frozen=True)
class TextExtractResult:
    ok: bool
    message: str
    text: str = ""


_RAPIDOCR = None


def _has_rapidocr() -> bool:
    return importlib.util.find_spec("rapidocr_onnxruntime") is not None


def _get_rapidocr():
    global _RAPIDOCR
    if _RAPIDOCR is None:
        from rapidocr_onnxruntime import RapidOCR

        _RAPIDOCR = RapidOCR()
    return _RAPIDOCR


def _rapidocr_to_text(ocr_res) -> str:
    if not ocr_res:
        return ""
    if isinstance(ocr_res, tuple) and ocr_res:
        ocr_res = ocr_res[0]
    if not ocr_res:
        return ""

    def sort_key(item):
        box = item[0] if isinstance(item, (list, tuple)) and len(item) > 0 else None
        if isinstance(box, (list, tuple)) and box:
            xs = [p[0] for p in box if isinstance(p, (list, tuple)) and len(p) >= 2]
            ys = [p[1] for p in box if isinstance(p, (list, tuple)) and len(p) >= 2]
            if xs and ys:
                return (min(ys), min(xs))
        return (0, 0)

    lines: list[str] = []
    for item in sorted(ocr_res, key=sort_key):
        if isinstance(item, (list, tuple)) and len(item) >= 2:
            text = str(item[1]).strip()
            if text:
                lines.append(text)
    return "\n".join(lines).strip()


def _ocr_image_rapidocr(image_input) -> str:
    if not _has_rapidocr():
        raise RuntimeError("RapidOCR is not available. Install rapidocr-onnxruntime.")
    engine = _get_rapidocr()
    try:
        res = engine(image_input)
    except Exception as e:
        raise RuntimeError(f"RapidOCR failed: {e}") from e
    return _rapidocr_to_text(res)


def _ocr_image_ai(file_path: Path, *, ai_manager, lang: str | None = None) -> str:
    try:
        if hasattr(ai_manager, "ensure_kind"):
            ai_manager.ensure_kind("vision")
    except Exception:
        pass
    if not ai_manager or not getattr(ai_manager, "is_ready", False):
        raise RuntimeError("AI model is not loaded.")
    if not getattr(ai_manager, "is_vision_model", False):
        raise RuntimeError("AI model does not support vision.")
    if not getattr(ai_manager, "model", None):
        raise RuntimeError("AI model is not available.")

    prep = getattr(ai_manager, "_prepare_image", None)
    if not prep:
        raise RuntimeError("AI image preparation is not available.")
    image_data = prep(Path(file_path))
    if not image_data:
        return ""

    lang_hint = f" Language hint: {lang}." if lang else ""
    prompt = (
        "Extract all readable text from this image. Preserve line breaks when possible."
        f"{lang_hint} If no text is visible, return an empty string."
    )
    response = ai_manager.create_chat_completion_safe(
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_data}"}},
                ],
            }
        ],
        temperature=0.1,
        max_tokens=512,
    )
    text = (response["choices"][0]["message"]["content"] or "").strip()
    # Strip code fences if present.
    text = re.sub(r"^```.*?\\n|```$", "", text, flags=re.DOTALL).strip()
    return text


def _ocr_image_ai_from_pil(image, *, ai_manager, lang: str | None = None) -> str:
    try:
        if hasattr(ai_manager, "ensure_kind"):
            ai_manager.ensure_kind("vision")
    except Exception:
        pass
    if not ai_manager or not getattr(ai_manager, "is_ready", False):
        raise RuntimeError("AI model is not loaded.")
    if not getattr(ai_manager, "is_vision_model", False):
        raise RuntimeError("AI model does not support vision.")
    if not getattr(ai_manager, "model", None):
        raise RuntimeError("AI model is not available.")

    buf = io.BytesIO()
    image.save(buf, format="PNG")
    data = base64.b64encode(buf.getvalue()).decode("ascii")
    if not data:
        return ""

    lang_hint = f" Language hint: {lang}." if lang else ""
    prompt = (
        "Extract all readable text from this image. Preserve line breaks when possible."
        f"{lang_hint} If no text is visible, return an empty string."
    )
    response = ai_manager.create_chat_completion_safe(
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{data}"}},
                ],
            }
        ],
        temperature=0.1,
        max_tokens=512,
    )
    text = (response["choices"][0]["message"]["content"] or "").strip()
    text = re.sub(r"^```.*?\\n|```$", "", text, flags=re.DOTALL).strip()
    return text


def _ocr_image(
    file_path: Path,
    *,
    engine: str,
    lang: str,
    ai_manager=None,
) -> str:
    engine = (engine or "auto").strip().lower()
    use_ai = engine == "ai"
    use_rapid = engine == "rapidocr"
    if engine == "auto":
        use_rapid = _has_rapidocr()
        use_ai = not use_rapid

    if use_rapid:
        from PIL import Image

        img = Image.open(file_path)
        if img.mode != "RGB":
            img = img.convert("RGB")
        import numpy as np

        return _ocr_image_rapidocr(np.array(img))

    if use_ai:
        return _ocr_image_ai(file_path, ai_manager=ai_manager, lang=lang)

    raise RuntimeError("No OCR engine available.")


def _ocr_pdf(
    pdf_path: Path,
    *,
    engine: str,
    lang: str,
    max_pages: int,
    ai_manager=None,
    progress_cb: Optional[Callable[[int, int], None]] = None,
) -> str:
    engine = (engine or "auto").strip().lower()
    if engine == "auto":
        engine = "rapidocr" if _has_rapidocr() else "ai"
    if engine == "rapidocr" and not _has_rapidocr():
        raise RuntimeError("RapidOCR is not available. Install rapidocr-onnxruntime.")
    try:
        import fitz  # type: ignore
    except Exception as e:
        raise RuntimeError("PDF OCR requires PyMuPDF (pymupdf).") from e

    from PIL import Image

    text_parts: list[str] = []
    with fitz.open(pdf_path) as doc:
        total = doc.page_count
        limit = total if int(max_pages) <= 0 else min(total, int(max_pages))
        for idx in range(limit):
            if progress_cb:
                progress_cb(idx + 1, limit)
            page = doc.load_page(idx)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            if engine == "ai":
                page_text = _ocr_image_ai_from_pil(img, ai_manager=ai_manager, lang=lang)
            else:
                import numpy as np

                page_text = _ocr_image_rapidocr(np.array(img))
            if page_text:
                text_parts.append(page_text)
    return "\n\n".join(text_parts).strip()


def extract_text_from_file(
    file_path: Path,
    *,
    ocr_mode: str = "auto",
    ocr_engine: str = "auto",
    ocr_lang: str = "eng",
    max_ocr_pages: int = 5,
    ai_manager=None,
    progress_cb: Optional[Callable[[int, int], None]] = None,
    max_chars: int = 500_000,
) -> TextExtractResult:
    """
    Extract text from a supported file.
    ocr_mode: auto | force | text_only
    ocr_engine: auto | rapidocr | ai
    """
    path = Path(file_path)
    if not path.exists() or not path.is_file():
        return TextExtractResult(ok=False, message="File not found.")

    ext = path.suffix.lower()
    ocr_mode = (ocr_mode or "auto").strip().lower()

    try:
        if ext in TEXT_EXTS:
            text = path.read_text(encoding="utf-8", errors="ignore")
            return TextExtractResult(ok=True, message="Text loaded.", text=text[:max_chars])

        if ext == ".pdf":
            from core.pdf_text_extract import extract_pdf_text

            text = extract_pdf_text(path)
            needs_ocr = (ocr_mode == "force") or (ocr_mode == "auto" and not (text or "").strip())
            if needs_ocr:
                try:
                    ocr_text = _ocr_pdf(
                        path,
                        engine=ocr_engine,
                        lang=ocr_lang,
                        max_pages=max_ocr_pages,
                        ai_manager=ai_manager,
                        progress_cb=progress_cb,
                    )
                    if ocr_text:
                        text = (text + "\n\n" + ocr_text).strip() if text else ocr_text
                except Exception as e:
                    if not (text or "").strip():
                        return TextExtractResult(ok=False, message=str(e))
            return TextExtractResult(ok=True, message="PDF extracted.", text=(text or "")[:max_chars])

        if ext in IMAGE_EXTS:
            if ocr_mode == "text_only":
                return TextExtractResult(ok=True, message="Skipped OCR (text-only).", text="")
            text = _ocr_image(path, engine=ocr_engine, lang=ocr_lang, ai_manager=ai_manager)
            return TextExtractResult(ok=True, message="OCR complete.", text=(text or "")[:max_chars])

        if ext in DOC_EXTS:
            try:
                import docx
            except Exception as e:
                return TextExtractResult(ok=False, message="DOCX extraction requires python-docx.")
            doc = docx.Document(path)
            parts = [p.text for p in doc.paragraphs if p.text]
            text = "\n".join(parts).strip()
            return TextExtractResult(ok=True, message="DOCX extracted.", text=text[:max_chars])

        if ext in PPT_EXTS:
            try:
                import pptx
            except Exception:
                return TextExtractResult(ok=False, message="PPTX extraction requires python-pptx.")
            prs = pptx.Presentation(path)
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    text = getattr(shape, "text", None)
                    if text:
                        parts.append(text)
            text_out = "\n".join(parts).strip()
            return TextExtractResult(ok=True, message="PPTX extracted.", text=text_out[:max_chars])

        if ext in SHEET_EXTS:
            try:
                import openpyxl
            except Exception:
                return TextExtractResult(ok=False, message="XLSX extraction requires openpyxl.")
            wb = openpyxl.load_workbook(path, data_only=True)
            parts: list[str] = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    vals = [str(v) for v in row if v is not None]
                    if vals:
                        parts.append("\t".join(vals))
            text_out = "\n".join(parts).strip()
            return TextExtractResult(ok=True, message="XLSX extracted.", text=text_out[:max_chars])

        return TextExtractResult(ok=False, message=f"Unsupported file type: {ext}")

    except Exception as e:
        return TextExtractResult(ok=False, message=str(e))


def normalize_text(text: str) -> str:
    text = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    # Remove excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Trim trailing spaces
    text = "\n".join(line.rstrip() for line in text.splitlines())
    return text.strip()
